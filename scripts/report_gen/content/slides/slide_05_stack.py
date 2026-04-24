"""Slide 5 — Technology stack (three columns)."""

from pptx.util import Inches, Emu
from pptx.enum.shapes import MSO_SHAPE

from ...pptx_utils import (
    bg,
    textbox,
    panel,
    accent_dot,
    slide_header,
    footer,
    solid,
    no_line,
    TEXT,
    TEXT_DIM,
    ACCENT,
    PANEL,
    BORDER,
)
from ._common import add_blank_slide
from ..project_facts import BACKEND_STACK, FRONTEND_STACK, INFRA_STACK


# Only the first 7 rows of each stack are shown to keep the column dense
# but readable. Values not shown here are present in the report tables.
_BACKEND_SLIDE = [(n, d) for n, _v, d in BACKEND_STACK[:7]]
_FRONTEND_SLIDE = [(n, d) for n, _v, d in FRONTEND_STACK[:7]]
_INFRA_SLIDE = list(INFRA_STACK[:7])


def _render_column(slide, x, y, width, height, title, items):
    """Draw a single stack column."""
    panel(slide, x, y, width, height, fill=PANEL)

    # Top accent rule.
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      x, y, width, Inches(0.08))
    solid(top_bar, ACCENT); no_line(top_bar)

    textbox(slide, Emu(x + Inches(0.3)), Emu(y + Inches(0.25)),
            width, Inches(0.5), title, size=20, bold=True, color=TEXT)

    # Divider under the title.
    sep = slide.shapes.add_connector(1,
                                      Emu(x + Inches(0.3)),
                                      Emu(y + Inches(0.85)),
                                      Emu(x + width - Inches(0.3)),
                                      Emu(y + Inches(0.85)))
    sep.line.color.rgb = BORDER
    from pptx.util import Pt as _Pt
    sep.line.width = _Pt(0.5)

    for j, (name, desc) in enumerate(items):
        row_y = Emu(y + Inches(1.05 + j * 0.5))
        accent_dot(slide, Emu(x + Inches(0.3)),
                   Emu(row_y + Inches(0.08)), size=Inches(0.1))
        textbox(slide, Emu(x + Inches(0.55)), row_y,
                Inches(3.5), Inches(0.3),
                name, size=12, bold=True, color=TEXT)
        textbox(slide, Emu(x + Inches(0.55)), Emu(row_y + Inches(0.22)),
                Inches(3.3), Inches(0.3),
                desc, size=10, color=TEXT_DIM)


def render(prs):
    """Render the three-column stack slide."""
    slide = add_blank_slide(prs)
    bg(slide)
    slide_header(slide, "Stack", "Tecnologias escolhidas",
                 "Ferramentas modernas, versões recentes, sem "
                 "deprecations.")

    col_w = Inches(4.0); col_h = Inches(4.6); col_y = Inches(2.4)
    columns = [
        (Inches(0.6), "Backend", _BACKEND_SLIDE),
        (Inches(4.75), "Frontend", _FRONTEND_SLIDE),
        (Inches(8.9), "Infraestrutura", _INFRA_SLIDE),
    ]
    for x, title, items in columns:
        _render_column(slide, x, col_y, col_w, col_h, title, items)

    footer(slide, 5)

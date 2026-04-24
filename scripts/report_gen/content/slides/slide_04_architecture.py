"""Slide 4 — System architecture."""

from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

from ...pptx_utils import (
    bg,
    textbox,
    panel,
    chip,
    slide_header,
    footer,
    solid,
    no_line,
    TEXT,
    TEXT_DIM,
    TEXT_MUTED,
    ACCENT,
    PANEL,
    SLIDE_W,
)
from ._common import add_blank_slide


def render(prs):
    """Render the four-block architecture diagram with arrows."""
    slide = add_blank_slide(prs)
    bg(slide)
    slide_header(slide, "Arquitetura", "Camadas bem separadas",
                 "SPA React  →  API FastAPI stateless  →  PostgreSQL "
                 "gerido.")

    # Horizontal flow diagram — 4 uniformly sized cards with arrows.
    blocks = [
        ("Browser", "Client", "React 19 SPA"),
        ("Vercel", "Edge CDN", "Build + HTTPS"),
        ("Render", "App Host", "FastAPI + Uvicorn"),
        ("PostgreSQL", "Database", "SQLAlchemy 2.0"),
    ]
    bw = Inches(2.7); bh = Inches(2.0); by = Inches(2.55)
    total_w = bw * 4 + Inches(0.4) * 3
    bx0 = Emu((SLIDE_W - total_w) / 2)

    for i, (title, sub1, sub2) in enumerate(blocks):
        left = Emu(bx0 + (bw + Inches(0.4)) * i)
        panel(slide, left, by, bw, bh, fill=PANEL)

        # Top accent rule on every card — visual rhythm device.
        top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                          left, by, bw, Inches(0.08))
        solid(top_bar, ACCENT); no_line(top_bar)

        chip(slide, Emu(left + Inches(0.25)),
             Emu(by + Inches(0.2)), sub1, fg=ACCENT, size=9)
        textbox(slide, Emu(left + Inches(0.25)),
                Emu(by + Inches(0.5)),
                bw, Inches(0.6), title, size=22, bold=True, color=TEXT)
        textbox(slide, Emu(left + Inches(0.25)),
                Emu(by + Inches(1.15)),
                Inches(2.4), Inches(0.4), sub2, size=11,
                color=TEXT_DIM, font="Consolas")

        # Connector arrow to the next card.
        if i < 3:
            ax1 = Emu(left + bw + Inches(0.05))
            ax2 = Emu(left + bw + Inches(0.35))
            ay = Emu(by + bh / 2)
            conn = slide.shapes.add_connector(1, ax1, ay, ax2, ay)
            conn.line.color.rgb = ACCENT
            conn.line.width = Pt(2.5)

    # Protocol/annotation row under the diagram.
    textbox(slide, Inches(0.6), Inches(4.8), Inches(12.1), Inches(0.4),
            "HTTP/JSON  ·  JWT stateless 7d  ·  CORS restrito  ·  "
            "OpenAPI auto em /docs",
            size=12, color=TEXT_MUTED, align=PP_ALIGN.CENTER,
            font="Consolas")

    # Two highlight panels: backend / frontend.
    hl_y = Inches(5.4)
    panel(slide, Inches(0.6), hl_y, Inches(6.0), Inches(1.5), fill=PANEL)
    chip(slide, Inches(0.8), Emu(hl_y + Inches(0.2)),
         "Backend", fg=ACCENT, size=10)
    textbox(slide, Inches(0.8), Emu(hl_y + Inches(0.5)),
            Inches(5.7), Inches(0.4),
            "19 routers  ·  100+ endpoints",
            size=15, bold=True, color=TEXT)
    textbox(slide, Inches(0.8), Emu(hl_y + Inches(0.9)),
            Inches(5.7), Inches(0.6),
            "Lifespan, migrations leves, OpenAPI, modo IA opcional com "
            "fallback de regras.",
            size=11, color=TEXT_DIM)

    panel(slide, Inches(6.75), hl_y, Inches(5.95), Inches(1.5), fill=PANEL)
    chip(slide, Inches(6.95), Emu(hl_y + Inches(0.2)),
         "Frontend", fg=ACCENT, size=10)
    textbox(slide, Inches(6.95), Emu(hl_y + Inches(0.5)),
            Inches(5.7), Inches(0.4),
            "17 páginas  ·  13 componentes",
            size=15, bold=True, color=TEXT)
    textbox(slide, Inches(6.95), Emu(hl_y + Inches(0.9)),
            Inches(5.7), Inches(0.6),
            "Context API, routing protegido por JWT, service layer, "
            "dark/light mode completo.",
            size=11, color=TEXT_DIM)

    footer(slide, 4)

"""
Shared helpers used by multiple slide modules.

Kept deliberately small — anything that's used by only one slide lives
inside that slide's module instead.
"""

from pptx.util import Inches, Emu
from pptx.enum.shapes import MSO_SHAPE

from ...pptx_utils import solid, no_line


def add_blank_slide(prs):
    """
    Append a blank slide (layout 6) to the presentation.

    :param prs: python-pptx Presentation.
    :returns: The new slide.
    """
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)


def accent_bar(slide, left, top, width=Inches(0.08),
               height=Inches(0.5), rgb=None):
    """
    Short solid bar used as a vertical or horizontal accent element.

    :param slide: python-pptx slide.
    :param left: Left offset.
    :param top: Top offset.
    :param width: Bar width.
    :param height: Bar height.
    :param rgb: Fill colour (RGBColor).
    :returns: The bar shape.
    """
    # Import here to avoid a circular import in pptx_utils package init.
    from ...pptx_utils import ACCENT
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top,
                                  width, height)
    solid(bar, rgb if rgb is not None else ACCENT)
    no_line(bar)
    return bar

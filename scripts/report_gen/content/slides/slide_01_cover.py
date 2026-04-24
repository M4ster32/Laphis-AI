"""Slide 1 — Cover."""

from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE

from ...pptx_utils import (
    bg,
    textbox,
    panel,
    chip,
    gradient_fill,
    solid,
    no_line,
    ACCENT,
    ACCENT_2,
    ACCENT_SOFT,
    BG_DEEP,
    TEXT,
    TEXT_DIM,
    PANEL,
    BORDER,
    SLIDE_H,
)
from ..project_facts import PROJECT
from ._common import add_blank_slide


def render(prs):
    """Render the cover slide with gradient stripe and identity blocks."""
    slide = add_blank_slide(prs)
    bg(slide)

    # Vertical gradient stripe on the far left (blue → violet).
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     0, 0, Inches(0.35), SLIDE_H)
    gradient_fill(stripe, ACCENT, ACCENT_2, angle=90)
    no_line(stripe)

    # Soft radial-feeling decoration using a large gradient oval.
    deco = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                   Inches(9.5), Inches(-2),
                                   Inches(7), Inches(7))
    gradient_fill(deco, ACCENT_SOFT, BG_DEEP, angle=135)
    no_line(deco)

    # University chip.
    chip(slide, Inches(1.0), Inches(0.9),
         f"{PROJECT['university']} · {PROJECT['faculty']}",
         fg=ACCENT, size=11)

    # Wordmark.
    textbox(slide, Inches(1.0), Inches(2.4), Inches(11), Inches(1.5),
            PROJECT["name"], size=96, bold=True, color=TEXT)

    # Full name — mixed styling to emphasise "Intelligent System".
    tb = slide.shapes.add_textbox(Inches(1.0), Inches(3.9),
                                   Inches(11), Inches(0.6))
    tf = tb.text_frame
    tf.margin_left = Emu(0)
    tf.margin_top = Emu(0)
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = "Life and Physical Health "
    r1.font.name = "Calibri"; r1.font.size = Pt(20)
    r1.font.color.rgb = TEXT_DIM
    r2 = p.add_run()
    r2.text = "Intelligent System"
    r2.font.name = "Calibri"; r2.font.size = Pt(20); r2.font.bold = True
    r2.font.color.rgb = ACCENT

    # Document title and pitch.
    textbox(slide, Inches(1.0), Inches(4.8), Inches(11), Inches(0.6),
            "Apresentação de Projeto II", size=22, bold=True, color=TEXT)
    textbox(slide, Inches(1.0), Inches(5.35), Inches(11), Inches(0.5),
            PROJECT["tagline"] + " — do tracking diário à IA "
            "conversacional.",
            size=13, color=TEXT_DIM, italic=True)

    # Footer metadata cards (group / course / date).
    footer_y = Inches(6.4)
    mini_w = Inches(3.0); mini_h = Inches(0.7)
    cards = [
        ("GRUPO", PROJECT["group"]),
        ("CURSO", "Engenharia Informática"),
        ("DATA", PROJECT["date"]),
    ]
    for i, (label, value) in enumerate(cards):
        left = Inches(1.0 + i * 3.4)
        panel(slide, left, footer_y, mini_w, mini_h,
              fill=PANEL, border=BORDER)
        textbox(slide, Emu(left + Inches(0.2)),
                Emu(footer_y + Inches(0.1)),
                mini_w, Inches(0.25),
                label, size=9, bold=True, color=ACCENT)
        textbox(slide, Emu(left + Inches(0.2)),
                Emu(footer_y + Inches(0.32)),
                mini_w, Inches(0.35),
                value, size=13, bold=True, color=TEXT)

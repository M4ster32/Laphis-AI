"""
Mid-level slide primitives: textboxes, panels, chips, background fills.

These compose the low-level fills helpers into reusable building blocks
that slide modules can use declaratively.
"""

from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from .theme import (
    BG_DEEP,
    BG_GRAD_END,
    PANEL,
    BORDER,
    TEXT,
    ACCENT,
    SLIDE_W,
    SLIDE_H,
)
from .fills import gradient_fill, solid, no_line, set_line


def bg(slide):
    """
    Fill the whole slide with a subtle diagonal gradient.

    :param slide: python-pptx slide.
    :returns: The inserted rectangle shape.
    """
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    no_line(rect)
    gradient_fill(rect, BG_DEEP, BG_GRAD_END, angle=135)
    return rect


def textbox(slide, left, top, width, height, text, size=14, bold=False,
            color=TEXT, align=PP_ALIGN.LEFT, font="Calibri",
            anchor=MSO_ANCHOR.TOP, italic=False):
    """
    Insert a single-paragraph textbox with the given run properties.

    :returns: The textbox shape.
    """
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    # Zero margins to match design mockups (python-pptx defaults are large).
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def multi_text(slide, left, top, width, height, lines,
               align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """
    Insert a textbox with multiple styled paragraphs.

    :param lines: Iterable of (text, size, bold, color[, font, space_after]).
    :returns: The textbox shape.
    """
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    for i, item in enumerate(lines):
        text, size, bold, color = item[0], item[1], item[2], item[3]
        font = item[4] if len(item) > 4 else "Calibri"
        space = item[5] if len(item) > 5 else 2
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space)
        run = p.add_run()
        run.text = text
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def panel(slide, left, top, width, height, fill=PANEL, border=BORDER,
          corner=0.06, border_w=0.75):
    """
    Rounded rectangle used as a generic card background.

    :returns: The rounded rectangle shape.
    """
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top,
                                   width, height)
    shape.adjustments[0] = corner
    solid(shape, fill)
    shape.line.color.rgb = border
    shape.line.width = Pt(border_w)
    return shape


def gradient_panel(slide, left, top, width, height, c_from, c_to,
                   angle=135, corner=0.06):
    """
    Rounded rectangle with a gradient fill and no outline.

    :returns: The rounded rectangle shape.
    """
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top,
                                   width, height)
    shape.adjustments[0] = corner
    gradient_fill(shape, c_from, c_to, angle=angle)
    no_line(shape)
    return shape


def accent_dot(slide, left, top, size=Inches(0.2), rgb=ACCENT):
    """
    Small solid circle used as a bullet marker.

    :returns: The oval shape.
    """
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    solid(dot, rgb)
    no_line(dot)
    return dot


def chip(slide, left, top, text, fg=ACCENT, size=10, width=None):
    """
    Uppercase, bold, coloured mini-label used as a kicker above titles.

    :param width: Optional width in EMU. Defaults to 2 inches, but long
        kicker labels (e.g. the cover slide university line) need more
        room or they wrap and the renderer hides the first line.
    :returns: The textbox.
    """
    tb = slide.shapes.add_textbox(
        left, top, width if width is not None else Inches(2), Inches(0.3)
    )
    tf = tb.text_frame
    tf.word_wrap = False
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text.upper()
    run.font.name = "Calibri"
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = fg
    return tb

"""
Gera apresentação LAPHIS em PowerPoint (.pptx) — 12 slides.
Versão premium: gradientes, cards polidos, KPIs fortes, dados reais do projeto.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn, nsmap
from lxml import etree

OUTPUT = "docs/LAPHIS_Apresentacao.pptx"

# ---------------------------------------------------------------------------
# PALETA — Premium Monochromatic (dark) com acento azul elétrico
# ---------------------------------------------------------------------------
BG_DEEP = RGBColor(0x05, 0x07, 0x0C)       # fundo principal
BG_GRAD_END = RGBColor(0x0D, 0x11, 0x1C)   # fundo gradiente final
PANEL = RGBColor(0x10, 0x14, 0x1F)         # painel base
PANEL_HI = RGBColor(0x16, 0x1B, 0x28)      # painel elevado
BORDER = RGBColor(0x23, 0x2A, 0x3D)        # bordas
BORDER_HI = RGBColor(0x33, 0x3D, 0x57)     # borda destacada
TEXT = RGBColor(0xF5, 0xF7, 0xFA)          # texto principal
TEXT_DIM = RGBColor(0xA5, 0xAE, 0xC0)      # texto secundário
TEXT_MUTED = RGBColor(0x6B, 0x74, 0x88)    # texto sutil
ACCENT = RGBColor(0x4F, 0x8C, 0xF7)        # azul elétrico
ACCENT_2 = RGBColor(0x8B, 0x5C, 0xF6)      # violeta (gradiente)
ACCENT_SOFT = RGBColor(0x26, 0x3D, 0x80)   # azul escuro (painéis de destaque)
SUCCESS = RGBColor(0x22, 0xC5, 0x5E)
WARN = RGBColor(0xF5, 0x9E, 0x0B)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ---------------------------------------------------------------------------
# HELPERS
# ---------------------------------------------------------------------------

def _nofill_line(shape):
    shape.line.fill.background()


def solid(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb


def gradient_fill(shape, rgb_from, rgb_to, angle=90):
    """Gradient linear via XML direto (python-pptx não tem API high-level)."""
    spPr = shape.fill._xPr
    # remover fill atual
    for tag in ("a:solidFill", "a:gradFill", "a:noFill", "a:blipFill", "a:pattFill"):
        for el in spPr.findall(qn(tag)):
            spPr.remove(el)
    grad = etree.SubElement(spPr, qn("a:gradFill"))
    grad.set("flip", "none")
    grad.set("rotWithShape", "1")
    gsLst = etree.SubElement(grad, qn("a:gsLst"))
    for pos, color in ((0, rgb_from), (100000, rgb_to)):
        gs = etree.SubElement(gsLst, qn("a:gs"))
        gs.set("pos", str(pos))
        srgb = etree.SubElement(gs, qn("a:srgbClr"))
        srgb.set("val", f"{color[0]:02X}{color[1]:02X}{color[2]:02X}"
                 if isinstance(color, tuple) else str(color))
    lin = etree.SubElement(grad, qn("a:lin"))
    lin.set("ang", str(angle * 60000))
    lin.set("scaled", "0")
    etree.SubElement(grad, qn("a:tileRect"))


def set_line(shape, rgb, width=0.75):
    shape.line.color.rgb = rgb
    shape.line.width = Pt(width)


def bg(slide):
    """Fundo com gradiente subtil."""
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    _nofill_line(rect)
    # converter RGBColor para tuple
    c1 = (BG_DEEP[0], BG_DEEP[1], BG_DEEP[2])
    c2 = (BG_GRAD_END[0], BG_GRAD_END[1], BG_GRAD_END[2])
    gradient_fill(rect, c1, c2, angle=135)


def textbox(slide, left, top, width, height, text, size=14, bold=False,
            color=TEXT, align=PP_ALIGN.LEFT, font="Calibri",
            anchor=MSO_ANCHOR.TOP, italic=False, spacing=None):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    if spacing:
        p.space_after = Pt(spacing)
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def multi_text(slide, left, top, width, height, lines, align=PP_ALIGN.LEFT,
               anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    for i, item in enumerate(lines):
        text, size, bold, color = item[0], item[1], item[2], item[3]
        font = item[4] if len(item) > 4 else "Calibri"
        space = item[5] if len(item) > 5 else 2
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space)
        run = p.add_run()
        run.text = text
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def panel(slide, left, top, width, height, fill=PANEL, border=BORDER,
          corner=0.06, border_w=0.75):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top,
                                   width, height)
    shape.adjustments[0] = corner
    solid(shape, fill)
    shape.line.color.rgb = border
    shape.line.width = Pt(border_w)
    return shape


def gradient_panel(slide, left, top, width, height, c_from, c_to, angle=135,
                    corner=0.06):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top,
                                   width, height)
    shape.adjustments[0] = corner
    gradient_fill(shape, (c_from[0], c_from[1], c_from[2]),
                  (c_to[0], c_to[1], c_to[2]), angle=angle)
    _nofill_line(shape)
    return shape


def accent_dot(slide, left, top, size=Inches(0.2), rgb=ACCENT):
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    solid(dot, rgb)
    _nofill_line(dot)
    return dot


def chip(slide, left, top, text, fg=ACCENT, bg_color=None, size=10):
    """Pequena pill/chip com texto."""
    tb = slide.shapes.add_textbox(left, top, Inches(2), Inches(0.3))
    tf = tb.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text.upper()
    run.font.name = "Calibri"
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = fg
    return tb


def slide_header(slide, kicker, title, subtitle=None):
    """Cabeçalho consistente em todos os slides internos."""
    # kicker colorido
    chip(slide, Inches(0.6), Inches(0.55), kicker, fg=ACCENT, size=11)
    # título grande
    textbox(slide, Inches(0.6), Inches(0.9), Inches(12), Inches(0.95),
            title, size=34, bold=True, color=TEXT)
    if subtitle:
        textbox(slide, Inches(0.6), Inches(1.75), Inches(12), Inches(0.5),
                subtitle, size=14, color=TEXT_DIM, italic=True)
    # linha fina de separação
    y_line = Inches(2.15) if subtitle else Inches(1.95)
    line = slide.shapes.add_connector(1, Inches(0.6), y_line,
                                       Inches(12.7), y_line)
    line.line.color.rgb = BORDER
    line.line.width = Pt(0.5)


def footer(slide, num, total=12):
    # linha
    line = slide.shapes.add_connector(1, Inches(0.6), Inches(7.05),
                                       Inches(12.7), Inches(7.05))
    line.line.color.rgb = BORDER
    line.line.width = Pt(0.5)
    # marca à esquerda
    textbox(slide, Inches(0.6), Inches(7.15), Inches(6), Inches(0.25),
            "LAPHIS  ·  Projeto II  ·  2026", size=9, color=TEXT_MUTED, bold=True)
    # número à direita
    textbox(slide, Inches(11.5), Inches(7.15), Inches(1.2), Inches(0.25),
            f"{num:02d} / {total:02d}", size=9, color=TEXT_MUTED,
            align=PP_ALIGN.RIGHT, bold=True)


# ---------------------------------------------------------------------------
# PRESENTATION SETUP
# ---------------------------------------------------------------------------

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]

# ===========================================================================
# SLIDE 1 — CAPA
# ===========================================================================
s = prs.slides.add_slide(blank)
bg(s)

# Faixa vertical de gradiente à esquerda
stripe = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.35), SLIDE_H)
gradient_fill(stripe, (ACCENT[0], ACCENT[1], ACCENT[2]),
              (ACCENT_2[0], ACCENT_2[1], ACCENT_2[2]), angle=90)
_nofill_line(stripe)

# Elemento decorativo: círculo grande desfocado à direita
deco = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.5), Inches(-2),
                           Inches(7), Inches(7))
gradient_fill(deco, (ACCENT_SOFT[0], ACCENT_SOFT[1], ACCENT_SOFT[2]),
              (BG_DEEP[0], BG_DEEP[1], BG_DEEP[2]), angle=135)
_nofill_line(deco)

# TAG no topo
chip(s, Inches(1.0), Inches(0.9), "Universidade de Trás-os-Montes e Alto Douro",
     fg=ACCENT, size=11)

# Logo grande
textbox(s, Inches(1.0), Inches(2.4), Inches(11), Inches(1.5),
        "LAPHIS", size=96, bold=True, color=TEXT)

# Tagline com cor
tag_tb = s.shapes.add_textbox(Inches(1.0), Inches(3.9), Inches(11), Inches(0.6))
tf = tag_tb.text_frame
tf.margin_left = Emu(0); tf.margin_top = Emu(0)
p = tf.paragraphs[0]
run = p.add_run()
run.text = "Life and Physical Health "
run.font.name = "Calibri"; run.font.size = Pt(20); run.font.color.rgb = TEXT_DIM
run2 = p.add_run()
run2.text = "Intelligent System"
run2.font.name = "Calibri"; run2.font.size = Pt(20); run2.font.bold = True
run2.font.color.rgb = ACCENT

# Título do documento
textbox(s, Inches(1.0), Inches(4.8), Inches(11), Inches(0.6),
        "Apresentação de Projeto II", size=22, bold=True, color=TEXT)
textbox(s, Inches(1.0), Inches(5.35), Inches(11), Inches(0.5),
        "Plataforma inteligente de saúde e fitness — do tracking diário à IA conversacional.",
        size=13, color=TEXT_DIM, italic=True)

# Rodapé grupo / data — cards mini
footer_y = Inches(6.4)
mini_w = Inches(3.0); mini_h = Inches(0.7)

for i, (label, value) in enumerate([
    ("GRUPO", "M4ster32"),
    ("CURSO", "Engenharia Informática"),
    ("DATA", "24 Abril 2026"),
]):
    left = Inches(1.0 + i * 3.4)
    panel(s, left, footer_y, mini_w, mini_h, fill=PANEL, border=BORDER)
    textbox(s, Emu(left + Inches(0.2)), Emu(footer_y + Inches(0.1)),
            mini_w, Inches(0.25), label, size=9, bold=True, color=ACCENT)
    textbox(s, Emu(left + Inches(0.2)), Emu(footer_y + Inches(0.32)),
            mini_w, Inches(0.35), value, size=13, bold=True, color=TEXT)

# ===========================================================================
# SLIDE 2 — O PROBLEMA
# ===========================================================================
s = prs.slides.add_slide(blank)
bg(s)
slide_header(s, "Contexto", "Um mercado fragmentado",
             "Tracking de saúde hoje: muitas apps, pouca inteligência cruzada.")

# Lead paragraph
textbox(s, Inches(0.6), Inches(2.4), Inches(12.1), Inches(0.9),
        "O utilizador típico usa 4-5 apps separadas. Nenhuma vê o todo — e as recomendações "
        "que oferecem são tabelas genéricas, não ajustadas ao perfil real.",
        size=15, color=TEXT_DIM)

# 3 cards de problema, com KPIs decorativos
card_w = Inches(4.0)
card_h = Inches(3.2)
card_y = Inches(3.6)
x0 = Inches(0.6)
gap = Inches(0.15)

problems = [
    ("01", "Silos de dados",
     "Treino no Strava. Comida no MyFitnessPal. Sono no Fitbit. Os dados nunca se cruzam — e é no cruzamento que está a inteligência."),
    ("02", "Recomendações genéricas",
     "Planos tabelados ignoram objetivos, restrições alimentares, nível de atividade real e feedback pessoal. Um plano igual para toda a gente."),
    ("03", "Zero adaptação",
     "Se o utilizador estagna, ninguém ajusta. Se progride, ninguém reage. O plano de dia 1 é o plano de dia 100."),
]

for i, (num, title, desc) in enumerate(problems):
    left = Emu(x0 + (card_w + gap) * i)
    card = panel(s, left, card_y, card_w, card_h, fill=PANEL)
    # número grande decorativo
    textbox(s, Emu(left + Inches(0.25)), Emu(card_y + Inches(0.2)),
            Inches(1.5), Inches(0.6), num, size=28, bold=True, color=ACCENT_SOFT)
    # título
    textbox(s, Emu(left + Inches(0.25)), Emu(card_y + Inches(0.85)),
            Inches(3.6), Inches(0.5), title, size=20, bold=True, color=TEXT)
    # barra de acento
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Emu(left + Inches(0.25)),
                              Emu(card_y + Inches(1.4)),
                              Inches(1.2), Inches(0.04))
    solid(bar, ACCENT); _nofill_line(bar)
    # desc
    textbox(s, Emu(left + Inches(0.25)), Emu(card_y + Inches(1.65)),
            Inches(3.55), Inches(1.5), desc, size=11.5, color=TEXT_DIM)

footer(s, 2)

# ===========================================================================
# SLIDE 3 — SOLUÇÃO
# ===========================================================================
s = prs.slides.add_slide(blank)
bg(s)
slide_header(s, "Proposta", "Uma plataforma. Todos os domínios.",
             "LAPHIS integra tracking + recomendação + IA conversacional num só produto.")

# HERO — painel grande com gradiente
hero_y = Inches(2.5)
hero = gradient_panel(s, Inches(0.6), hero_y, Inches(12.1), Inches(1.8),
                       ACCENT_SOFT, PANEL_HI, angle=90)
textbox(s, Inches(1.0), Emu(hero_y + Inches(0.3)), Inches(11), Inches(0.55),
        "Plataforma web com IA que cruza todos os domínios num único perfil.",
        size=22, bold=True, color=TEXT)
# Pilares inline
pillars = ["Treino", "Nutrição", "Hidratação", "Peso", "Meditação", "Plano Diário"]
px0 = Inches(1.0)
pgap = Inches(0.1)
pw = Inches(1.73)
py = Emu(hero_y + Inches(1.05))
for i, p_name in enumerate(pillars):
    left = Emu(px0 + (pw + pgap) * i)
    p_panel = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, py, pw, Inches(0.5))
    p_panel.adjustments[0] = 0.4
    solid(p_panel, BG_DEEP); set_line(p_panel, BORDER_HI, 0.5)
    textbox(s, left, Emu(py + Inches(0.13)), pw, Inches(0.3),
            p_name, size=11, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

# Objetivos — numerados
textbox(s, Inches(0.6), Inches(4.7), Inches(12), Inches(0.45),
        "Objetivos de Projeto II", size=16, bold=True, color=TEXT)
textbox(s, Inches(0.6), Inches(5.1), Inches(12), Inches(0.35),
        "Do protótipo funcional ao produto em produção.", size=11,
        color=TEXT_MUTED, italic=True)

objs = [
    ("01", "Plataforma web completa"),
    ("02", "Motor de recomendação"),
    ("03", "Assistente IA + RAG"),
    ("04", "Design system premium"),
    ("05", "Testes & deploy real"),
]
ow = Inches(2.35); oh = Inches(1.3); oy = Inches(5.6); ogap = Inches(0.125)
ox = Inches(0.6)
for i, (num, text) in enumerate(objs):
    left = Emu(ox + (ow + ogap) * i)
    p_card = panel(s, left, oy, ow, oh, fill=PANEL)
    textbox(s, Emu(left + Inches(0.25)), Emu(oy + Inches(0.15)),
            ow, Inches(0.4), num, size=16, bold=True, color=ACCENT)
    textbox(s, Emu(left + Inches(0.25)), Emu(oy + Inches(0.55)),
            Inches(2.1), Inches(0.7), text, size=12, bold=True, color=TEXT)

footer(s, 3)

# ===========================================================================
# SLIDE 4 — ARQUITETURA
# ===========================================================================
s = prs.slides.add_slide(blank)
bg(s)
slide_header(s, "Arquitetura", "Camadas bem separadas",
             "SPA React  →  API FastAPI stateless  →  PostgreSQL gerido.")

# Diagrama grande horizontal
blocks = [
    ("Browser", "Client", "React 19 SPA"),
    ("Vercel", "Edge CDN", "Build + HTTPS"),
    ("Render", "App Host", "FastAPI + Uvicorn"),
    ("PostgreSQL", "Database", "SQLAlchemy 2.0"),
]
bw = Inches(2.7); bh = Inches(2.0)
by = Inches(2.55)
total_w = bw * 4 + Inches(0.4) * 3
bx0 = Emu((SLIDE_W - total_w) / 2)

for i, (title, sub1, sub2) in enumerate(blocks):
    left = Emu(bx0 + (bw + Inches(0.4)) * i)
    # painel
    pnl = panel(s, left, by, bw, bh, fill=PANEL, border=BORDER)
    # barra de topo colorida
    top_bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, by, bw, Inches(0.08))
    solid(top_bar, ACCENT); _nofill_line(top_bar)
    # kicker
    chip(s, Emu(left + Inches(0.25)), Emu(by + Inches(0.2)),
         sub1, fg=ACCENT, size=9)
    # title
    textbox(s, Emu(left + Inches(0.25)), Emu(by + Inches(0.5)),
            bw, Inches(0.6), title, size=22, bold=True, color=TEXT)
    # sub
    textbox(s, Emu(left + Inches(0.25)), Emu(by + Inches(1.15)),
            Inches(2.4), Inches(0.4), sub2, size=11, color=TEXT_DIM,
            font="Consolas")
    # seta
    if i < 3:
        ax1 = Emu(left + bw + Inches(0.05))
        ax2 = Emu(left + bw + Inches(0.35))
        ay = Emu(by + bh / 2)
        conn = s.shapes.add_connector(1, ax1, ay, ax2, ay)
        conn.line.color.rgb = ACCENT
        conn.line.width = Pt(2.5)

# Legenda do fluxo
textbox(s, Inches(0.6), Inches(4.8), Inches(12.1), Inches(0.4),
        "HTTP/JSON  ·  JWT stateless 7d  ·  CORS restrito  ·  OpenAPI auto em /docs",
        size=12, color=TEXT_MUTED, align=PP_ALIGN.CENTER, font="Consolas")

# Destaques backend/frontend
hl_y = Inches(5.4)
panel(s, Inches(0.6), hl_y, Inches(6.0), Inches(1.5), fill=PANEL)
chip(s, Inches(0.8), Emu(hl_y + Inches(0.2)), "Backend", fg=ACCENT, size=10)
textbox(s, Inches(0.8), Emu(hl_y + Inches(0.5)), Inches(5.7), Inches(0.4),
        "19 routers  ·  100+ endpoints", size=15, bold=True, color=TEXT)
textbox(s, Inches(0.8), Emu(hl_y + Inches(0.9)), Inches(5.7), Inches(0.6),
        "Lifespan, migrations leves, OpenAPI, modo IA opcional com fallback de regras.",
        size=11, color=TEXT_DIM)

panel(s, Inches(6.75), hl_y, Inches(5.95), Inches(1.5), fill=PANEL)
chip(s, Inches(6.95), Emu(hl_y + Inches(0.2)), "Frontend", fg=ACCENT, size=10)
textbox(s, Inches(6.95), Emu(hl_y + Inches(0.5)), Inches(5.7), Inches(0.4),
        "17 páginas  ·  13 componentes", size=15, bold=True, color=TEXT)
textbox(s, Inches(6.95), Emu(hl_y + Inches(0.9)), Inches(5.7), Inches(0.6),
        "Context API, routing protegido por JWT, service layer, dark/light mode completo.",
        size=11, color=TEXT_DIM)

footer(s, 4)

# ===========================================================================
# SLIDE 5 — STACK
# ===========================================================================
s = prs.slides.add_slide(blank)
bg(s)
slide_header(s, "Stack", "Tecnologias escolhidas",
             "Ferramentas modernas, versões recentes, sem deprecations.")

col_w = Inches(4.0)
col_h = Inches(4.6)
col_y = Inches(2.4)
col_x = [Inches(0.6), Inches(4.75), Inches(8.9)]
titles = ["Backend", "Frontend", "Infraestrutura"]

stacks = [
    [
        ("Python 3.11+", "Linguagem base"),
        ("FastAPI", "Async · OpenAPI auto"),
        ("SQLAlchemy 2.0", "ORM moderno"),
        ("Pydantic 2", "Validação type-safe"),
        ("PyJWT + bcrypt", "Auth segura"),
        ("OpenAI SDK", "Chat · planos · RAG"),
        ("PyMuPDF", "Extração PDF (RAG)"),
    ],
    [
        ("React 19", "Concurrent UI"),
        ("Vite", "HMR instantâneo"),
        ("React Router 7", "Routing aninhado"),
        ("Recharts", "Gráficos SVG"),
        ("Lucide React", "Ícones consistentes"),
        ("jsPDF", "Export relatórios"),
        ("CSS Custom Props", "Theming"),
    ],
    [
        ("Vercel", "Hosting frontend"),
        ("Render", "Backend + PostgreSQL"),
        ("GitHub", "CI/CD automático"),
        ("pytest + httpx", "Testes integrados"),
        ("ESLint", "Lint JS/React"),
        ("Docker", "Build reprodutível"),
        ("SMTP", "Emails transacionais"),
    ],
]

for i in range(3):
    # painel com barra superior
    panel(s, col_x[i], col_y, col_w, col_h, fill=PANEL)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, col_x[i], col_y,
                              col_w, Inches(0.08))
    solid(bar, ACCENT); _nofill_line(bar)
    # título
    textbox(s, Emu(col_x[i] + Inches(0.3)), Emu(col_y + Inches(0.25)),
            col_w, Inches(0.5), titles[i], size=20, bold=True, color=TEXT)
    # separator
    line = s.shapes.add_connector(1,
                                   Emu(col_x[i] + Inches(0.3)),
                                   Emu(col_y + Inches(0.85)),
                                   Emu(col_x[i] + col_w - Inches(0.3)),
                                   Emu(col_y + Inches(0.85)))
    line.line.color.rgb = BORDER; line.line.width = Pt(0.5)
    # items
    for j, (name, desc) in enumerate(stacks[i]):
        y = Emu(col_y + Inches(1.05 + j * 0.5))
        # dot
        accent_dot(s, Emu(col_x[i] + Inches(0.3)), Emu(y + Inches(0.08)),
                   size=Inches(0.1))
        textbox(s, Emu(col_x[i] + Inches(0.55)), y, Inches(3.5), Inches(0.3),
                name, size=12, bold=True, color=TEXT)
        textbox(s, Emu(col_x[i] + Inches(0.55)), Emu(y + Inches(0.22)),
                Inches(3.3), Inches(0.3), desc, size=10, color=TEXT_DIM)

footer(s, 5)

# ===========================================================================
# SLIDE 6 — MODELO DE DADOS
# ===========================================================================
s = prs.slides.add_slide(blank)
bg(s)
slide_header(s, "Dados", "Modelo de domínio",
             "17 entidades ligadas ao utilizador. Datetimes timezone-aware, JSON flexível.")

# User central com halo
cx = Inches(6.3); cy = Inches(4.2)
halo = s.shapes.add_shape(MSO_SHAPE.OVAL, Emu(cx - Inches(0.35)),
                           Emu(cy - Inches(0.35)), Inches(2.7), Inches(1.7))
solid(halo, ACCENT_SOFT); _nofill_line(halo)

user_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, cy,
                                Inches(2.0), Inches(1.0))
user_box.adjustments[0] = 0.25
gradient_fill(user_box, (ACCENT[0], ACCENT[1], ACCENT[2]),
              (ACCENT_2[0], ACCENT_2[1], ACCENT_2[2]), angle=135)
_nofill_line(user_box)
textbox(s, cx, Emu(cy + Inches(0.2)), Inches(2.0), Inches(0.3),
        "USER", size=10, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
textbox(s, cx, Emu(cy + Inches(0.45)), Inches(2.0), Inches(0.4),
        "entidade raiz", size=11, color=RGBColor(0xE0, 0xE7, 0xFF),
        align=PP_ALIGN.CENTER, italic=True)

# Entidades à volta
entities = [
    # nome, x, y, categoria (cor)
    ("Profile", Inches(2.3), Inches(2.5), "core"),
    ("Category", Inches(5.2), Inches(2.5), "core"),
    ("ProgressSnapshot", Inches(8.0), Inches(2.5), "core"),
    ("Plan", Inches(10.5), Inches(3.0), "ai"),
    ("DailyPlan", Inches(10.8), Inches(4.2), "ai"),
    ("AdaptationLog", Inches(10.5), Inches(5.4), "ai"),
    ("ChatSession", Inches(9.3), Inches(6.0), "ai"),
    ("ChatMessage", Inches(6.8), Inches(6.3), "ai"),
    ("ZenSession", Inches(4.3), Inches(6.3), "log"),
    ("WorkoutLog", Inches(2.0), Inches(6.0), "log"),
    ("MealLog", Inches(0.8), Inches(5.1), "log"),
    ("WaterLog", Inches(0.8), Inches(3.9), "log"),
    ("WeightEntry", Inches(0.8), Inches(2.8), "log"),
]

cat_colors = {"core": ACCENT, "ai": ACCENT_2, "log": SUCCESS}

# Ligações primeiro
for name, ex, ey, cat in entities:
    conn = s.shapes.add_connector(1,
                                   Emu(cx + Inches(1.0)), Emu(cy + Inches(0.5)),
                                   Emu(ex + Inches(0.8)), Emu(ey + Inches(0.2)))
    conn.line.color.rgb = BORDER
    conn.line.width = Pt(0.5)

# Caixas das entidades
for name, ex, ey, cat in entities:
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, ex, ey,
                              Inches(1.65), Inches(0.45))
    box.adjustments[0] = 0.3
    solid(box, PANEL)
    box.line.color.rgb = cat_colors[cat]
    box.line.width = Pt(1.25)
    textbox(s, ex, Emu(ey + Inches(0.08)), Inches(1.65), Inches(0.3),
            name, size=10, bold=True, color=TEXT, align=PP_ALIGN.CENTER)

# Legenda
leg_y = Inches(6.55)
leg_items = [("Core", ACCENT), ("IA / Planos", ACCENT_2), ("Tracking", SUCCESS)]
lx = Inches(0.6)
for lbl, color in leg_items:
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, lx, Emu(leg_y + Inches(0.08)),
                              Inches(0.2), Inches(0.2))
    solid(dot, color); _nofill_line(dot)
    textbox(s, Emu(lx + Inches(0.3)), leg_y, Inches(1.8), Inches(0.3),
            lbl, size=11, color=TEXT_DIM)
    lx = Emu(lx + Inches(2.0))

footer(s, 6)

# ===========================================================================
# SLIDE 7 — BACKEND
# ===========================================================================
s = prs.slides.add_slide(blank)
bg(s)
slide_header(s, "Backend", "19 routers, arquitetura modular",
             "Cada domínio num router próprio. Separação clara entre API e lógica de negócio.")

# KPI bar no topo
kpi_y = Inches(2.4)
kpis = [
    ("19", "routers"),
    ("100+", "endpoints"),
    ("17", "entidades"),
    ("857", "linhas no recommender"),
    ("~7k", "linhas backend"),
]
kw = Inches(2.4); kh = Inches(1.0); kgap = Inches(0.08)
for i, (val, lbl) in enumerate(kpis):
    left = Emu(Inches(0.6) + (kw + kgap) * i)
    panel(s, left, kpi_y, kw, kh, fill=PANEL_HI)
    textbox(s, left, Emu(kpi_y + Inches(0.1)), kw, Inches(0.55),
            val, size=26, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    textbox(s, left, Emu(kpi_y + Inches(0.65)), kw, Inches(0.3),
            lbl, size=10, color=TEXT_DIM, align=PP_ALIGN.CENTER)

# Grid de routers
grid_y = Inches(3.7)
routers = [
    ("auth", 14), ("profile", 4), ("logs", 6), ("plans", 11), ("chat", 8),
    ("daily_plan", 6), ("adaptation", 11), ("progress", 8), ("reports", 7),
    ("water", 5), ("weight", 5), ("zen", 5), ("categories", 5),
    ("exercises", 6), ("rag_ingest", 4), ("rag_ask", 4), ("ingest", 1),
    ("ask", 3), ("health", 1),
]
cols = 5
cw = Inches(2.4); ch = Inches(0.6); gx = Inches(0.1); gy = Inches(0.1)
for i, (name, count) in enumerate(routers):
    row = i // cols; col = i % cols
    left = Emu(Inches(0.6) + (cw + gx) * col)
    top = Emu(grid_y + (ch + gy) * row)
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, cw, ch)
    box.adjustments[0] = 0.2
    solid(box, PANEL); set_line(box, BORDER, 0.5)
    # route name
    textbox(s, Emu(left + Inches(0.2)), Emu(top + Inches(0.15)),
            Inches(1.6), Inches(0.3), f"/{name}", size=11, bold=True,
            color=TEXT, font="Consolas")
    # endpoint count chip
    textbox(s, Emu(left + cw - Inches(0.55)), Emu(top + Inches(0.15)),
            Inches(0.4), Inches(0.3), str(count), size=11, bold=True,
            color=ACCENT, align=PP_ALIGN.RIGHT)

footer(s, 7)

# ===========================================================================
# SLIDE 8 — FRONTEND
# ===========================================================================
s = prs.slides.add_slide(blank)
bg(s)
slide_header(s, "Frontend", "Experiência polida do login ao plano diário",
             "17 páginas, 13 componentes, ~15k linhas de JSX/CSS.")

# Hero card grande à esquerda com mockup simbólico
hero_y = Inches(2.4)
panel(s, Inches(0.6), hero_y, Inches(6.2), Inches(4.5), fill=PANEL)
# barra de topo simulando janela de app
for i, color in enumerate([WARN, ACCENT, SUCCESS]):
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL,
                              Emu(Inches(0.85) + Inches(0.3) * i),
                              Emu(hero_y + Inches(0.25)),
                              Inches(0.2), Inches(0.2))
    solid(dot, color); _nofill_line(dot)

# "URL bar"
url = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                          Inches(2.0), Emu(hero_y + Inches(0.22)),
                          Inches(4.6), Inches(0.3))
url.adjustments[0] = 0.4
solid(url, BG_DEEP); set_line(url, BORDER, 0.5)
textbox(s, Inches(2.2), Emu(hero_y + Inches(0.25)), Inches(4.4), Inches(0.25),
        "laphis.vercel.app/app/dashboard", size=9,
        color=TEXT_DIM, font="Consolas")

# Conteúdo do mockup
textbox(s, Inches(0.85), Emu(hero_y + Inches(0.8)), Inches(5.8), Inches(0.35),
        "Bom dia, Henrique", size=16, bold=True, color=TEXT)
textbox(s, Inches(0.85), Emu(hero_y + Inches(1.15)), Inches(5.8), Inches(0.3),
        "Plano de hoje — terça, 24 abril", size=10,
        color=TEXT_DIM, italic=True)

# 4 metric boxes no mockup
metrics = [("1 820", "kcal"), ("142 g", "proteína"),
           ("2.1 L", "água"), ("72.4 kg", "peso")]
mw = Inches(1.3); mh = Inches(0.85); my = Emu(hero_y + Inches(1.65))
for i, (val, lbl) in enumerate(metrics):
    left = Emu(Inches(0.85) + (mw + Inches(0.1)) * i)
    mpan = panel(s, left, my, mw, mh, fill=BG_DEEP, border=BORDER)
    textbox(s, left, Emu(my + Inches(0.1)), mw, Inches(0.4),
            val, size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    textbox(s, left, Emu(my + Inches(0.5)), mw, Inches(0.3),
            lbl, size=9, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

# gráfico simbólico (bars)
gy = Emu(hero_y + Inches(2.7))
textbox(s, Inches(0.85), gy, Inches(5.8), Inches(0.3),
        "Semana", size=10, bold=True, color=TEXT_DIM)
bar_vals = [0.5, 0.7, 0.4, 0.85, 0.6, 0.9, 0.75]
bar_w = Inches(0.5); bgap = Inches(0.25); bmax_h = Inches(1.0)
for i, v in enumerate(bar_vals):
    h = Emu(bmax_h * v)
    left = Emu(Inches(0.85) + (bar_w + bgap) * i)
    top = Emu(gy + Inches(0.35) + (bmax_h - h))
    bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, bar_w, h)
    bar.adjustments[0] = 0.3
    solid(bar, ACCENT); _nofill_line(bar)

# Coluna da direita — lista de páginas
rx = Inches(7.1)
textbox(s, rx, Inches(2.4), Inches(5.6), Inches(0.4),
        "Páginas principais", size=14, bold=True, color=ACCENT)
pages = [
    ("Dashboard", "KPIs diários, shortcuts, plano do dia"),
    ("Logs", "treinos e refeições num só fluxo"),
    ("Plans · PlanDetail", "planos gerados e editáveis"),
    ("DailyPlan", "plano adaptativo com feedback"),
    ("Chat", "assistente IA com contexto do perfil"),
    ("Reports", "gráficos semanais + export PDF"),
    ("Zen", "meditação guiada"),
    ("Profile · Settings", "perfil biométrico + preferências"),
]
for i, (name, desc) in enumerate(pages):
    y = Inches(2.85 + i * 0.5)
    accent_dot(s, rx, Emu(y + Inches(0.08)), size=Inches(0.12))
    textbox(s, Emu(rx + Inches(0.2)), y, Inches(2.2), Inches(0.3),
            name, size=11.5, bold=True, color=TEXT)
    textbox(s, Emu(rx + Inches(2.5)), y, Inches(3.2), Inches(0.3),
            desc, size=10, color=TEXT_DIM)

footer(s, 8)

# ===========================================================================
# SLIDE 9 — IA & RECOMENDAÇÃO
# ===========================================================================
s = prs.slides.add_slide(blank)
bg(s)
slide_header(s, "Inteligência", "Motor próprio + IA opcional",
             "Funciona sem API key (modo regras). Com chave, ativa chat + RAG + planos por LLM.")

# 2 colunas grandes
col_y = Inches(2.4)
col_h = Inches(4.6)

# ESQUERDA — Regras
left_col = Inches(0.6)
panel(s, left_col, col_y, Inches(6.0), col_h, fill=PANEL)
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_col, col_y,
                          Inches(6.0), Inches(0.08))
solid(bar, SUCCESS); _nofill_line(bar)
chip(s, Emu(left_col + Inches(0.3)), Emu(col_y + Inches(0.25)),
     "Sempre ativo", fg=SUCCESS, size=10)
textbox(s, Emu(left_col + Inches(0.3)), Emu(col_y + Inches(0.55)),
        Inches(5.5), Inches(0.5), "Motor de Regras", size=22,
        bold=True, color=TEXT)
textbox(s, Emu(left_col + Inches(0.3)), Emu(col_y + Inches(1.05)),
        Inches(5.5), Inches(0.35),
        "core/recommender.py  ·  857 linhas Python puro",
        size=10, color=TEXT_MUTED, font="Consolas")

rules_data = [
    ("BMI", "peso / altura²"),
    ("TDEE", "Harris-Benedict × atividade"),
    ("Macros", "split personalizado por objetivo"),
    ("Hidratação", "30-35 ml/kg de peso"),
    ("Plano diário", "treinos + refeições + macros"),
    ("Adaptação", "AdaptationLog regista alterações"),
]
for i, (k, v) in enumerate(rules_data):
    y = Emu(col_y + Inches(1.6 + i * 0.47))
    # k como "tag"
    tag_w = Inches(1.6)
    tag = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Emu(left_col + Inches(0.3)), y, tag_w, Inches(0.35))
    tag.adjustments[0] = 0.4
    solid(tag, BG_DEEP); set_line(tag, BORDER, 0.5)
    textbox(s, Emu(left_col + Inches(0.3)), Emu(y + Inches(0.06)),
            tag_w, Inches(0.25), k, size=10, bold=True, color=ACCENT,
            align=PP_ALIGN.CENTER)
    textbox(s, Emu(left_col + Inches(2.0)), Emu(y + Inches(0.05)),
            Inches(4), Inches(0.3), v, size=11, color=TEXT_DIM)

# DIREITA — IA
right_col = Inches(6.75)
panel(s, right_col, col_y, Inches(5.95), col_h, fill=PANEL)
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, right_col, col_y,
                          Inches(5.95), Inches(0.08))
solid(bar, ACCENT); _nofill_line(bar)
chip(s, Emu(right_col + Inches(0.3)), Emu(col_y + Inches(0.25)),
     "Opcional · OpenAI", fg=ACCENT, size=10)
textbox(s, Emu(right_col + Inches(0.3)), Emu(col_y + Inches(0.55)),
        Inches(5.5), Inches(0.5), "Camada IA + RAG", size=22,
        bold=True, color=TEXT)
textbox(s, Emu(right_col + Inches(0.3)), Emu(col_y + Inches(1.05)),
        Inches(5.5), Inches(0.35),
        "core/ai_engine.py  ·  core/rag.py",
        size=10, color=TEXT_MUTED, font="Consolas")

ai_items = [
    ("Chat conversacional", "sessões persistentes com contexto do perfil"),
    ("Planos gerados", "prompt com objetivos + restrições + histórico"),
    ("Pipeline RAG", "PDF → PyMuPDF → chunks → resposta contextualizada"),
    ("Detecção automática", "sem chave, cai em modo regras sem crash"),
    ("Feedback loop", "PlanFeedback melhora prompts futuros"),
]
for i, (k, v) in enumerate(ai_items):
    y = Emu(col_y + Inches(1.6 + i * 0.56))
    accent_dot(s, Emu(right_col + Inches(0.3)), Emu(y + Inches(0.08)),
               size=Inches(0.14), rgb=ACCENT)
    textbox(s, Emu(right_col + Inches(0.55)), y, Inches(5), Inches(0.3),
            k, size=12, bold=True, color=TEXT)
    textbox(s, Emu(right_col + Inches(0.55)), Emu(y + Inches(0.26)),
            Inches(5.2), Inches(0.3), v, size=10, color=TEXT_DIM)

footer(s, 9)

# ===========================================================================
# SLIDE 10 — QUALIDADE
# ===========================================================================
s = prs.slides.add_slide(blank)
bg(s)
slide_header(s, "Qualidade", "Código sem dívida técnica",
             "62 testes, 0 warnings, todas as deprecations resolvidas.")

# 4 stat cards grandes
stats = [
    ("62", "testes", "pytest + httpx"),
    ("0", "warnings", "deprecations limpas"),
    ("7", "módulos", "cobertura dos domínios"),
    ("~20k", "linhas", "backend + frontend"),
]
sw = Inches(3.0); sh = Inches(1.9); sy = Inches(2.45); sgap = Inches(0.1)
sx0 = Inches(0.6)
for i, (val, lbl, sub) in enumerate(stats):
    left = Emu(sx0 + (sw + sgap) * i)
    # Painel com gradiente mais sutil para o card do "0"
    if val == "0":
        pnl = gradient_panel(s, left, sy, sw, sh, ACCENT_SOFT, PANEL_HI, angle=135)
    else:
        pnl = panel(s, left, sy, sw, sh, fill=PANEL)
    # barra superior
    top_bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, sy, sw, Inches(0.06))
    solid(top_bar, ACCENT); _nofill_line(top_bar)
    textbox(s, left, Emu(sy + Inches(0.35)), sw, Inches(0.9),
            val, size=52, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
    textbox(s, left, Emu(sy + Inches(1.15)), sw, Inches(0.3),
            lbl, size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    textbox(s, left, Emu(sy + Inches(1.45)), sw, Inches(0.3),
            sub, size=10, color=TEXT_DIM, align=PP_ALIGN.CENTER)

# Tabela de modernização
mod_y = Inches(4.7)
textbox(s, Inches(0.6), mod_y, Inches(12), Inches(0.4),
        "Modernização — deprecated → atual", size=15, bold=True, color=TEXT)
textbox(s, Inches(0.6), Emu(mod_y + Inches(0.38)), Inches(12), Inches(0.3),
        "Refactor completo para FastAPI / Pydantic v2 modernos.",
        size=11, color=TEXT_MUTED, italic=True)

mods = [
    ("@app.on_event('startup')", "lifespan(app) async context"),
    ("orm_mode / class Config", "model_config = ConfigDict(...)"),
    ("datetime.utcnow()", "datetime.now(timezone.utc)"),
    ("@validator (v1)", "@model_validator (v2)"),
]
for i, (old, new) in enumerate(mods):
    y = Emu(mod_y + Inches(0.85 + (i // 2) * 0.55))
    col_offset = (i % 2) * Inches(6.2)
    base_x = Inches(0.6) + col_offset
    # bloco old
    old_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   base_x, y, Inches(2.5), Inches(0.4))
    old_box.adjustments[0] = 0.3
    solid(old_box, PANEL); set_line(old_box, BORDER, 0.5)
    textbox(s, base_x, Emu(y + Inches(0.09)),
            Inches(2.5), Inches(0.25), old, size=9,
            color=TEXT_DIM, font="Consolas", align=PP_ALIGN.CENTER)
    # seta
    textbox(s, Emu(base_x + Inches(2.55)), Emu(y + Inches(0.05)),
            Inches(0.4), Inches(0.3), "→", size=14, bold=True,
            color=ACCENT, align=PP_ALIGN.CENTER)
    # bloco new
    new_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Emu(base_x + Inches(2.95)), y,
                                   Inches(3.0), Inches(0.4))
    new_box.adjustments[0] = 0.3
    solid(new_box, ACCENT_SOFT); _nofill_line(new_box)
    textbox(s, Emu(base_x + Inches(2.95)), Emu(y + Inches(0.09)),
            Inches(3.0), Inches(0.25), new, size=9, bold=True,
            color=TEXT, font="Consolas", align=PP_ALIGN.CENTER)

footer(s, 10)

# ===========================================================================
# SLIDE 11 — DEPLOY & PRODUÇÃO
# ===========================================================================
s = prs.slides.add_slide(blank)
bg(s)
slide_header(s, "Produção", "Deploy automático, em produção",
             "Push to main → live em 30 segundos.")

# Pipeline visual
pipe_y = Inches(2.5)
textbox(s, Inches(0.6), pipe_y, Inches(12), Inches(0.4),
        "Pipeline CI/CD", size=14, bold=True, color=ACCENT)

flow = [
    ("git push", PANEL, TEXT),
    ("GitHub", PANEL, TEXT),
    ("Build", PANEL_HI, TEXT),
    ("Deploy", ACCENT_SOFT, TEXT),
    ("Live", SUCCESS, TEXT),
]
fw = Inches(2.1); fh = Inches(0.85); fy = Inches(3.0); fgap = Inches(0.3)
total = fw * 5 + fgap * 4
fx0 = Emu((SLIDE_W - total) / 2)
for i, (lbl, bg_c, fg_c) in enumerate(flow):
    left = Emu(fx0 + (fw + fgap) * i)
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, fy, fw, fh)
    box.adjustments[0] = 0.3
    solid(box, bg_c); set_line(box, BORDER_HI, 1.0)
    textbox(s, left, Emu(fy + Inches(0.25)), fw, Inches(0.4),
            lbl, size=14, bold=True, color=fg_c, align=PP_ALIGN.CENTER)
    if i < 4:
        arrow_x1 = Emu(left + fw + Inches(0.02))
        arrow_x2 = Emu(left + fw + fgap - Inches(0.02))
        arrow_y = Emu(fy + fh / 2)
        conn = s.shapes.add_connector(1, arrow_x1, arrow_y, arrow_x2, arrow_y)
        conn.line.color.rgb = ACCENT; conn.line.width = Pt(2)

# 2 cards Vercel + Render
card_y = Inches(4.4)
card_w = Inches(6.0); card_h = Inches(2.5)

# Vercel
panel(s, Inches(0.6), card_y, card_w, card_h, fill=PANEL)
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), card_y,
                          card_w, Inches(0.08))
solid(bar, TEXT); _nofill_line(bar)
chip(s, Inches(0.85), Emu(card_y + Inches(0.2)), "Frontend", fg=TEXT_DIM, size=10)
textbox(s, Inches(0.85), Emu(card_y + Inches(0.5)), Inches(5.5), Inches(0.5),
        "Vercel", size=22, bold=True, color=TEXT)
vlines = [
    ("~30s", "build e deploy automáticos"),
    ("Global", "CDN edge em 100+ regiões"),
    ("Auto", "HTTPS e certificados geridos"),
    ("Preview", "URL único por pull request"),
]
for i, (k, v) in enumerate(vlines):
    y = Emu(card_y + Inches(1.05 + i * 0.32))
    textbox(s, Inches(0.85), y, Inches(0.9), Inches(0.3),
            k, size=11, bold=True, color=ACCENT)
    textbox(s, Inches(1.8), y, Inches(4.3), Inches(0.3),
            v, size=11, color=TEXT_DIM)

# Render
panel(s, Inches(6.75), card_y, Inches(5.95), card_h, fill=PANEL)
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.75), card_y,
                          Inches(5.95), Inches(0.08))
solid(bar, ACCENT); _nofill_line(bar)
chip(s, Inches(7.0), Emu(card_y + Inches(0.2)),
     "Backend + BD", fg=ACCENT, size=10)
textbox(s, Inches(7.0), Emu(card_y + Inches(0.5)), Inches(5.5), Inches(0.5),
        "Render", size=22, bold=True, color=TEXT)
rlines = [
    ("FastAPI", "Uvicorn + workers auto"),
    ("Postgres", "Base de dados gerida"),
    ("Health", "/health + auto-restart"),
    ("Logs", "dashboard em tempo real"),
]
for i, (k, v) in enumerate(rlines):
    y = Emu(card_y + Inches(1.05 + i * 0.32))
    textbox(s, Inches(7.0), y, Inches(0.9), Inches(0.3),
            k, size=11, bold=True, color=ACCENT)
    textbox(s, Inches(7.95), y, Inches(4.3), Inches(0.3),
            v, size=11, color=TEXT_DIM)

footer(s, 11)

# ===========================================================================
# SLIDE 12 — CONCLUSÃO
# ===========================================================================
s = prs.slides.add_slide(blank)
bg(s)

# Slide de fecho — mais visual, menos texto
# Fundo: grande gradiente decorativo
deco_big = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-3), Inches(-3),
                                Inches(12), Inches(12))
gradient_fill(deco_big, (ACCENT_SOFT[0], ACCENT_SOFT[1], ACCENT_SOFT[2]),
              (BG_DEEP[0], BG_DEEP[1], BG_DEEP[2]), angle=135)
_nofill_line(deco_big)

chip(s, Inches(0.6), Inches(0.55), "Conclusão", fg=ACCENT, size=11)
textbox(s, Inches(0.6), Inches(0.9), Inches(12), Inches(1.0),
        "Balanço & próximos passos", size=34, bold=True, color=TEXT)

# Linha de conquistas — KPIs grandes
kpi_y = Inches(2.3)
kpi_items = [
    ("12/12", "Requisitos funcionais"),
    ("8/8", "Requisitos não-funcionais"),
    ("100%", "Testes a passar"),
    ("LIVE", "Produção ativa"),
]
kw = Inches(3.0); kh = Inches(1.2); kgap = Inches(0.1)
for i, (val, lbl) in enumerate(kpi_items):
    left = Emu(Inches(0.6) + (kw + kgap) * i)
    pnl = panel(s, left, kpi_y, kw, kh, fill=PANEL_HI)
    top_bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, kpi_y, kw, Inches(0.06))
    solid(top_bar, SUCCESS); _nofill_line(top_bar)
    textbox(s, left, Emu(kpi_y + Inches(0.15)), kw, Inches(0.6),
            val, size=30, bold=True, color=SUCCESS, align=PP_ALIGN.CENTER)
    textbox(s, left, Emu(kpi_y + Inches(0.8)), kw, Inches(0.3),
            lbl, size=11, color=TEXT_DIM, align=PP_ALIGN.CENTER)

# 2 colunas finais — feito / futuro
col_y = Inches(3.9)
col_h = Inches(2.75)

# Feito
panel(s, Inches(0.6), col_y, Inches(6.0), col_h, fill=PANEL)
textbox(s, Inches(0.85), Emu(col_y + Inches(0.25)),
        Inches(5.5), Inches(0.4), "O que foi conseguido",
        size=16, bold=True, color=ACCENT)
wins = [
    "Aplicação em produção acessível a utilizadores reais",
    "Motor de recomendação próprio + IA opcional",
    "Pipeline RAG funcional sobre PDFs",
    "Design system coeso com dark/light mode completo",
    "Código modernizado, zero deprecations",
]
for i, w in enumerate(wins):
    y = Emu(col_y + Inches(0.75 + i * 0.38))
    check = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.85),
                                 Emu(y + Inches(0.05)),
                                 Inches(0.2), Inches(0.2))
    solid(check, SUCCESS); _nofill_line(check)
    textbox(s, Inches(1.2), y, Inches(5.3), Inches(0.3),
            w, size=11, color=TEXT)

# Futuro
panel(s, Inches(6.75), col_y, Inches(5.95), col_h, fill=PANEL)
textbox(s, Inches(7.0), Emu(col_y + Inches(0.25)),
        Inches(5.5), Inches(0.4), "Próximos passos",
        size=16, bold=True, color=ACCENT_2)
future = [
    "Notificações push e gamificação",
    "Integração com wearables (Apple Health, Fitbit)",
    "ML próprio treinado nos dados dos utilizadores",
    "App nativa (React Native) para iOS e Android",
    "Marketplace de planos validados por profissionais",
]
for i, f in enumerate(future):
    y = Emu(col_y + Inches(0.75 + i * 0.38))
    arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Inches(7.0),
                                 Emu(y + Inches(0.05)),
                                 Inches(0.2), Inches(0.2))
    solid(arrow, ACCENT_2); _nofill_line(arrow)
    textbox(s, Inches(7.35), y, Inches(5.3), Inches(0.3),
            f, size=11, color=TEXT)

# Slogan final
slogan_y = Inches(6.8)
textbox(s, Inches(0.6), slogan_y, Inches(12.1), Inches(0.4),
        "LAPHIS — saúde inteligente, em todos os domínios, numa única plataforma.",
        size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER, italic=True)
# pequeno obrigado
textbox(s, Inches(0.6), Emu(slogan_y + Inches(0.35)),
        Inches(12.1), Inches(0.3),
        "Obrigado.",
        size=11, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

# ---------------------------------------------------------------------------
# SAVE
# ---------------------------------------------------------------------------
prs.save(OUTPUT)
print(f"PowerPoint gerado: {OUTPUT}")
print(f"Slides: {len(prs.slides)}")
